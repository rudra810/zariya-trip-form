"""
Zariya Website Audit - Premium Gen Z Luxury PPT Generator
Creates a visually appealing, point-wise presentation
"""

from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

OUT_FILE = "Zariya_Website_Audit_Premium.pptx"

# Premium color palette
COLORS = {
    "black": RGBColor(18, 18, 18),
    "white": RGBColor(255, 255, 255),
    "cream": RGBColor(250, 248, 245),
    "gold": RGBColor(197, 165, 98),
    "gold_dark": RGBColor(165, 135, 75),
    "gray_dark": RGBColor(60, 60, 65),
    "gray_medium": RGBColor(120, 120, 128),
    "gray_light": RGBColor(200, 200, 205),
    "green": RGBColor(46, 125, 50),
    "red": RGBColor(198, 40, 40),
    "blue": RGBColor(30, 90, 150),
}


def set_slide_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_gold_accent_line(slide, left, top, width):
    """Add a gold accent line"""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["gold"]
    line.line.fill.background()


def add_title_slide(prs):
    """Slide 1: Cover"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, COLORS["black"])
    
    # Main title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "ZARIYA"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = COLORS["white"]
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.7), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "WEBSITE AUDIT REPORT"
    p.font.size = Pt(28)
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS["gold"]
    p.font.letter_spacing = Pt(4)
    
    # Gold accent line
    add_gold_accent_line(slide, 0.8, 4.3, 3.0)
    
    # Description
    desc_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.5))
    tf = desc_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Premium Gen Z Fashion Positioning Review"
    p.font.size = Pt(18)
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS["gray_light"]
    
    p2 = tf.add_paragraph()
    p2.text = f"March 2026"
    p2.font.size = Pt(16)
    p2.font.name = "Aptos"
    p2.font.color.rgb = COLORS["gray_medium"]
    
    # Bottom accent
    add_gold_accent_line(slide, 0.8, 6.8, 11.7)


def add_section_header(slide, title, subtitle="", dark_bg=False):
    """Add a consistent header to content slides"""
    title_color = COLORS["white"] if dark_bg else COLORS["black"]
    sub_color = COLORS["gray_light"] if dark_bg else COLORS["gray_medium"]
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.35), Inches(12.0), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title.upper()
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = title_color
    p.font.letter_spacing = Pt(2)
    
    # Gold underline
    add_gold_accent_line(slide, 0.7, 0.95, 2.5)
    
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.1), Inches(12.0), Inches(0.4))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(14)
        p.font.name = "Aptos"
        p.font.color.rgb = sub_color


def add_bullet_points(slide, items, left=0.7, top=1.5, width=12.0, dark_bg=False, with_icons=False):
    """Add formatted bullet points"""
    text_color = COLORS["gray_light"] if dark_bg else COLORS["gray_dark"]
    bullet_color = COLORS["gold"]
    
    current_top = top
    for item in items:
        # Gold bullet circle
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(current_top + 0.12), Inches(0.12), Inches(0.12))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = bullet_color
        bullet.line.fill.background()
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(left + 0.28), Inches(current_top), Inches(width - 0.5), Inches(0.6))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(16)
        p.font.name = "Aptos"
        p.font.color.rgb = text_color
        
        current_top += 0.7


def add_two_columns(slide, left_title, left_items, right_title, right_items, dark_bg=False):
    """Add two-column layout"""
    title_color = COLORS["white"] if dark_bg else COLORS["black"]
    text_color = COLORS["gray_light"] if dark_bg else COLORS["gray_dark"]
    
    # Left column title
    lt_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = lt_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = title_color
    
    # Left column line
    add_gold_accent_line(slide, 0.7, 1.85, 1.5)
    
    # Right column title
    rt_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.4), Inches(5.8), Inches(0.5))
    tf = rt_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = title_color
    
    # Right column line
    add_gold_accent_line(slide, 7.0, 1.85, 1.5)
    
    # Left items
    current_top = 2.1
    for item in left_items:
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.7), Inches(current_top + 0.1), Inches(0.1), Inches(0.1))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = COLORS["gold"]
        bullet.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(1.0), Inches(current_top), Inches(5.5), Inches(0.55))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.name = "Aptos"
        p.font.color.rgb = text_color
        current_top += 0.6
    
    # Right items
    current_top = 2.1
    for item in right_items:
        bullet = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.0), Inches(current_top + 0.1), Inches(0.1), Inches(0.1))
        bullet.fill.solid()
        bullet.fill.fore_color.rgb = COLORS["gold"]
        bullet.line.fill.background()
        
        text_box = slide.shapes.add_textbox(Inches(7.3), Inches(current_top), Inches(5.5), Inches(0.55))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(14)
        p.font.name = "Aptos"
        p.font.color.rgb = text_color
        current_top += 0.6


def add_numbered_list(slide, items, left=0.7, top=1.5, dark_bg=False):
    """Add numbered list with styled numbers"""
    text_color = COLORS["gray_light"] if dark_bg else COLORS["gray_dark"]
    
    current_top = top
    for i, item in enumerate(items, 1):
        # Number circle
        num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(current_top), Inches(0.35), Inches(0.35))
        num_circle.fill.solid()
        num_circle.fill.fore_color.rgb = COLORS["gold"]
        num_circle.line.fill.background()
        
        # Number text
        num_box = slide.shapes.add_textbox(Inches(left), Inches(current_top + 0.02), Inches(0.35), Inches(0.35))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i)
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS["black"]
        p.alignment = PP_ALIGN.CENTER
        
        # Text
        text_box = slide.shapes.add_textbox(Inches(left + 0.5), Inches(current_top + 0.02), Inches(11.5), Inches(0.55))
        tf = text_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(15)
        p.font.name = "Aptos"
        p.font.color.rgb = text_color
        
        current_top += 0.65


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ===== SLIDE 1: COVER =====
    add_title_slide(prs)
    
    # ===== SLIDE 2: AUDIT SCOPE =====
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s2, COLORS["cream"])
    add_section_header(s2, "Audit Scope", "What was analyzed in this review")
    add_bullet_points(s2, [
        "Homepage messaging, visual hierarchy, and conversion elements",
        "Collection page merchandising and pricing presentation",
        "Contact, shipping, and refund policy transparency",
        "Brand-market fit for premium Gen Z fashion segment",
        "Startup readiness: trust, retention, and scale factors",
        "Data source: Live website crawl (zariyaofficial.com)"
    ], top=1.6)
    
    # ===== SLIDE 3: EXECUTIVE SUMMARY =====
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s3, COLORS["black"])
    add_section_header(s3, "Executive Summary", "Key findings at a glance", dark_bg=True)
    add_bullet_points(s3, [
        "Strong positioning intent: Premium quality + personal storytelling",
        "Current site behaves like discount-led store, not luxury brand",
        "Urgent fixes needed: Brand credibility, product storytelling, trust proof",
        "Critical issue: About page appears broken (404 error)",
        "Policies exist but need premium-grade clarity and polish",
        "High potential for conversion & AOV improvement with focused redesign"
    ], top=1.6, dark_bg=True)
    
    # ===== SLIDE 4: WHAT'S WORKING =====
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s4, COLORS["cream"])
    add_section_header(s4, "What's Working Well", "Strengths to build upon")
    add_numbered_list(s4, [
        "Strong narrative: Quality, craftsmanship, and identity messaging",
        "Trust themes present: Quality assurance, safe checkout cues",
        "Direct social commerce: Instagram & WhatsApp integration",
        "Clear pricing visibility and simple catalog navigation",
        "Operational basics in place: Shipping & refund policies",
        "Mobile-friendly Shopify infrastructure for fast iteration"
    ], top=1.6)
    
    # ===== SLIDE 5: CRITICAL ISSUES =====
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s5, COLORS["cream"])
    add_section_header(s5, "Critical Issues", "Priority P0/P1 - Fix immediately")
    add_bullet_points(s5, [
        "Luxury vs Discount conflict: 41% OFF framing undermines premium image",
        "Brand proof deficit: Limited social proof, founder credibility, press mentions",
        "Broken credibility path: About page returns 404 error",
        "Product pages under-leveraged: Missing fabric, fit, process, care details",
        "AI-generated media names reduce authenticity perception",
        "Inconsistent language polish in policy content hurts trust"
    ], top=1.6)
    
    # ===== SLIDE 6: MINOR ISSUES =====
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s6, COLORS["cream"])
    add_section_header(s6, "Minor Issues", "Priority P2 - Address soon")
    add_bullet_points(s6, [
        "Homepage hero: Image-heavy with low proposition density",
        "Limited discoverability links in navigation structure",
        "Floating social widget may disrupt premium visual calm",
        "Policy copy contains placeholder-like formatting artifacts",
        "Missing shipping promise badges near purchase areas",
        "No loyalty/membership layer for customer retention"
    ], top=1.6)
    
    # ===== SLIDE 7: DESIGN PRINCIPLES =====
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s7, COLORS["cream"])
    add_section_header(s7, "Design Principles", "Current vs Recommended approach")
    add_two_columns(s7,
        "Current Pattern",
        [
            "Generic D2C grid layout",
            "Discount-first visual emphasis",
            "Limited craft storytelling",
            "Trust and style cues separated",
            "Low editorial depth"
        ],
        "Luxury Gen Z Pattern",
        [
            "Editorial commerce: Lookbook + product",
            "Material-first: GSM, weave, fit videos",
            "Curated UGC + style creator proof",
            "Coherent brand world aesthetic",
            "Controlled scarcity: Drops & waitlists"
        ]
    )
    
    # ===== SLIDE 8: STARTUP PRINCIPLES =====
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s8, COLORS["black"])
    add_section_header(s8, "Growth Principles", "Startup scaling strategies", dark_bg=True)
    add_bullet_points(s8, [
        "Trust Flywheel: Proof assets → Better conversion → More reviews → Lower CAC",
        "Retention Over Discounting: Reduce promo dependence, increase conviction",
        "LTV Architecture: Post-purchase journeys, drops calendar, loyalty tiers",
        "Merchandising Discipline: Hero SKUs, bundles, new-drop sequencing",
        "Measurement Stack: CVR by source, AOV by collection, repeat rate",
        "Founder Brand Signal: Narrative, behind-the-scenes, mission proof"
    ], top=1.6, dark_bg=True)
    
    # ===== SLIDE 9: COMPETITIVE ANALYSIS =====
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s9, COLORS["cream"])
    add_section_header(s9, "Competitive Position", "vs Premium D2C peers")
    add_two_columns(s9,
        "Where Zariya Leads",
        [
            "Clear quality + emotional identity",
            "Social messaging already integrated",
            "Simple, friction-light discovery",
            "Clean baseline storefront",
            "Quick iteration on Shopify"
        ],
        "Where Zariya Trails",
        [
            "Luxury authority systems",
            "Editorial storytelling depth",
            "Premium PDP conversion design",
            "Copy and policy polish",
            "Post-purchase retention"
        ]
    )
    
    # ===== SLIDE 10: 30-60-90 DAY PLAN =====
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s10, COLORS["cream"])
    add_section_header(s10, "Action Plan", "30-60-90 day roadmap")
    
    # Three columns for timeline
    phases = [
        ("Days 0-30", ["Fix About page", "Rewrite policy copy", "Add trust badges", "Improve hero message", "Redesign discount display"]),
        ("Days 31-60", ["Upgrade PDP content", "Add material/fit/care info", "Deploy UGC system", "Premium review moderation", "Visual consistency"]),
        ("Days 61-90", ["Launch drop calendar", "Build referral loop", "Founder content series", "KPI dashboard setup", "Weekly growth reviews"])
    ]
    
    for i, (phase_title, phase_items) in enumerate(phases):
        left = 0.7 + (i * 4.2)
        
        # Phase header
        header_box = s10.shapes.add_textbox(Inches(left), Inches(1.5), Inches(3.8), Inches(0.45))
        tf = header_box.text_frame
        p = tf.paragraphs[0]
        p.text = phase_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.name = "Aptos Display"
        p.font.color.rgb = COLORS["gold_dark"]
        
        add_gold_accent_line(s10, left, 1.9, 1.2)
        
        current_top = 2.1
        for item in phase_items:
            bullet = s10.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(current_top + 0.08), Inches(0.08), Inches(0.08))
            bullet.fill.solid()
            bullet.fill.fore_color.rgb = COLORS["gold"]
            bullet.line.fill.background()
            
            text_box = s10.shapes.add_textbox(Inches(left + 0.2), Inches(current_top), Inches(3.6), Inches(0.5))
            tf = text_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = item
            p.font.size = Pt(13)
            p.font.name = "Aptos"
            p.font.color.rgb = COLORS["gray_dark"]
            current_top += 0.5
    
    # ===== SLIDE 11: KPIs =====
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s11, COLORS["cream"])
    add_section_header(s11, "KPIs To Track", "Metrics that matter")
    add_two_columns(s11,
        "Conversion Quality",
        [
            "PDP to cart rate",
            "Checkout completion rate",
            "Average order value (AOV)",
            "Full-price vs discounted sales ratio",
            "Return rate by SKU"
        ],
        "Brand Health",
        [
            "Repeat purchase rate (30/60/90 day)",
            "NPS / Post-purchase satisfaction",
            "UGC volume & review velocity",
            "Direct traffic & branded search trend",
            "Contribution margin after returns"
        ]
    )
    
    # ===== SLIDE 12: EVIDENCE =====
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s12, COLORS["black"])
    add_section_header(s12, "Evidence Snapshot", "From website crawl analysis", dark_bg=True)
    add_bullet_points(s12, [
        "Homepage: Why Zariya, Quality Assurance, Brand Credibility, Safe Checkout",
        "Collection: Repeated 41% OFF, pricing Rs.699 vs Rs.1,199",
        "Contact: Instagram, WhatsApp, email, phone available",
        "Shipping: Process 1-3 days, delivery 3-6 business days",
        "Refund: 7-day return window, customer-paid return shipping",
        "About Page: Returns 404 - requires immediate correction"
    ], top=1.6, dark_bg=True)
    
    # ===== SLIDE 13: CLOSING =====
    s13 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s13, COLORS["black"])
    
    # Closing message
    close_box = s13.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(1.0))
    tf = close_box.text_frame
    p = tf.paragraphs[0]
    p.text = "THANK YOU"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.name = "Aptos Display"
    p.font.color.rgb = COLORS["white"]
    
    add_gold_accent_line(s13, 0.8, 3.8, 2.0)
    
    sub_box = s13.shapes.add_textbox(Inches(0.8), Inches(4.1), Inches(11.7), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Ready to elevate Zariya to premium status"
    p.font.size = Pt(20)
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS["gray_light"]
    
    # Footer
    footer = s13.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "zariyaofficial.com | March 2026"
    p.font.size = Pt(12)
    p.font.name = "Aptos"
    p.font.color.rgb = COLORS["gray_medium"]
    
    add_gold_accent_line(s13, 0.8, 6.8, 11.7)
    
    # Save
    prs.save(OUT_FILE)
    print(f"✓ Created: {OUT_FILE}")
    print(f"  - 13 slides with premium styling")
    print(f"  - Gold accents, clean typography")
    print(f"  - Point-wise formatting throughout")


if __name__ == "__main__":
    main()
