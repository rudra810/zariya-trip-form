from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT_FILE = "Zariya_Website_Audit_GenZ_Luxury_March2026.pptx"


# Simple visual system
COLORS = {
    "bg_dark": RGBColor(15, 16, 20),
    "bg_light": RGBColor(245, 243, 239),
    "title": RGBColor(255, 255, 255),
    "title_dark": RGBColor(22, 23, 28),
    "accent": RGBColor(201, 169, 97),
    "body_dark": RGBColor(50, 52, 58),
}


def style_title(shape, dark=False):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(36)
            r.font.bold = True
            r.font.name = "Aptos Display"
            r.font.color.rgb = COLORS["title_dark"] if dark else COLORS["title"]


def style_subtitle(shape, dark=False):
    tf = shape.text_frame
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(16)
            r.font.name = "Aptos"
            r.font.color.rgb = COLORS["body_dark"] if dark else RGBColor(225, 225, 230)


def set_slide_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_header_bar(slide, title, dark=False):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), Inches(12.0), Inches(0.6))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.name = "Aptos Display"
    p.font.color.rgb = COLORS["title_dark"] if dark else COLORS["title"]

    line = slide.shapes.add_shape(1, Inches(0.6), Inches(0.95), Inches(2.1), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = COLORS["accent"]
    line.line.fill.background()


def add_bullets(slide, items, left=0.8, top=1.4, width=12.0, height=5.8, dark=True):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.size = Pt(20)
        p.font.name = "Aptos"
        p.font.color.rgb = COLORS["body_dark"] if dark else RGBColor(236, 236, 240)


def add_two_column(slide, left_title, left_items, right_title, right_items):
    title_l = slide.shapes.add_textbox(Inches(0.8), Inches(1.35), Inches(5.8), Inches(0.5))
    tf_l = title_l.text_frame
    tf_l.text = left_title
    tf_l.paragraphs[0].font.bold = True
    tf_l.paragraphs[0].font.size = Pt(22)
    tf_l.paragraphs[0].font.name = "Aptos Display"
    tf_l.paragraphs[0].font.color.rgb = COLORS["title_dark"]

    title_r = slide.shapes.add_textbox(Inches(6.9), Inches(1.35), Inches(5.8), Inches(0.5))
    tf_r = title_r.text_frame
    tf_r.text = right_title
    tf_r.paragraphs[0].font.bold = True
    tf_r.paragraphs[0].font.size = Pt(22)
    tf_r.paragraphs[0].font.name = "Aptos Display"
    tf_r.paragraphs[0].font.color.rgb = COLORS["title_dark"]

    add_bullets(slide, left_items, left=0.8, top=1.9, width=5.8, height=5.0, dark=True)
    add_bullets(slide, right_items, left=6.9, top=1.9, width=5.8, height=5.0, dark=True)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Cover
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    set_slide_bg(s1, COLORS["bg_dark"])
    s1.shapes.title.text = "Zariya Website Audit"
    s1.placeholders[1].text = (
        "Luxury Gen Z Wear Positioning Review\n"
        "Design + Startup Principles\n"
        f"Date: {date.today().strftime('%d %b %Y')}"
    )
    style_title(s1.shapes.title)
    style_subtitle(s1.placeholders[1])

    accent = s1.shapes.add_shape(1, Inches(0.8), Inches(6.2), Inches(3.2), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = COLORS["accent"]
    accent.line.fill.background()

    # Slide 2: What was analyzed
    s2 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s2, COLORS["bg_light"])
    add_header_bar(s2, "Audit Scope", dark=True)
    add_bullets(
        s2,
        [
            "Homepage messaging, visual hierarchy, trust blocks, and conversion cues",
            "Collection page merchandising, discount signaling, and product-card UX",
            "Contact, shipping, and refund policy transparency",
            "Brand-market fit against premium Gen Z fashion expectations",
            "Startup readiness: trust, retention, repeat purchase, and scale risks",
            "Data source: live crawl of zariyaofficial.com pages on 29 Mar 2026",
        ],
        dark=True,
    )

    # Slide 3: Executive summary
    s3 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s3, COLORS["bg_light"])
    add_header_bar(s3, "Executive Summary", dark=True)
    add_bullets(
        s3,
        [
            "Positioning intent is strong: premium quality + personal storytelling is clear",
            "Current site behaves like a discount-led store, not a luxury-first brand",
            "Most urgent fixes: brand credibility architecture, product storytelling depth, and trust proof",
            "High impact issue: About page appears non-functional (404), reducing legitimacy",
            "Policies exist, but language and precision need premium-grade clarity",
            "With focused redesign and proof systems, conversion quality and AOV can improve significantly",
        ],
        dark=True,
    )

    # Slide 4: What is working
    s4 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s4, COLORS["bg_light"])
    add_header_bar(s4, "What Zariya Is Doing Well", dark=True)
    add_bullets(
        s4,
        [
            "Strong core narrative on homepage: quality, craftsmanship, and identity",
            "Relevant trust themes present: quality assurance, safe checkout, brand credibility",
            "Direct social commerce bridges via Instagram and WhatsApp",
            "Simple catalog flow and clear sale pricing visibility",
            "Operational basics are in place: shipping and return/refund policy pages",
            "Mobile-friendly Shopify infrastructure supports fast iteration",
        ],
        dark=True,
    )

    # Slide 5: Critical issues (major)
    s5 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s5, COLORS["bg_light"])
    add_header_bar(s5, "Major Gaps (Priority P0/P1)", dark=True)
    add_bullets(
        s5,
        [
            "Luxury vs discount conflict: prominent 41% OFF framing weakens premium perception",
            "Brand proof deficit: limited visible social proof, founder credibility, and press/customer proof",
            "Broken credibility path: About page resolves to 404 in crawl",
            "Product pages likely under-leveraged for luxury conversion (fabric, fit, process, care, craft depth)",
            "AI-generated style cues in media naming can reduce authenticity perception for premium audience",
            "Inconsistent language polish in policy content can reduce trust for high-intent buyers",
        ],
        dark=True,
    )

    # Slide 6: Minor issues
    s6 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s6, COLORS["bg_light"])
    add_header_bar(s6, "Minor Gaps (Priority P2)", dark=True)
    add_bullets(
        s6,
        [
            "Homepage hero appears image-heavy with low immediate proposition density",
            "Information architecture shows limited discoverability links in crawl output",
            "Floating social widget may compete with premium visual calm if overused",
            "Policy copy contains placeholder-like wording and formatting artifacts",
            "Lack of explicit shipping promise badges and easy returns summary near buy intent areas",
            "No clear loyalty, membership, or post-purchase relationship layer on core surfaces",
        ],
        dark=True,
    )

    # Slide 7: Design principles
    s7 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s7, COLORS["bg_light"])
    add_header_bar(s7, "Design Principles To Apply", dark=True)
    add_two_column(
        s7,
        "Current Pattern",
        [
            "Generic D2C grid rhythm",
            "Discount-first visual emphasis",
            "Limited craft storytelling blocks",
            "Trust and style cues separated",
            "Low editorial depth",
        ],
        "Recommended Luxury Gen Z Pattern",
        [
            "Editorial commerce: lookbook + product utility",
            "Material-first proof: GSM, weave, fit videos, care science",
            "Social proof with curation: verified UGC + style creators",
            "Brand world coherence: typography, spacing, photography direction",
            "Controlled scarcity: drops, waitlists, limited editions",
        ],
    )

    # Slide 8: Startup principles
    s8 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s8, COLORS["bg_light"])
    add_header_bar(s8, "Startup Growth Principles", dark=True)
    add_bullets(
        s8,
        [
            "Trust flywheel: proof assets -> stronger conversion -> better reviews -> lower CAC",
            "Retention over discounting: reduce promo dependence and increase product conviction",
            "LTV architecture: post-purchase journeys, drops calendar, loyalty tiers, referral loops",
            "Merchandising discipline: hero SKUs, bundles, and new-drop sequencing",
            "Measurement stack: CVR by source, AOV by collection, repeat rate, contribution margin",
            "Founder brand signal: narrative, behind-the-scenes quality process, and clear mission proof",
        ],
        dark=True,
    )

    # Slide 9: Competitive comparison
    s9 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s9, COLORS["bg_light"])
    add_header_bar(s9, "Compared To Premium D2C Peers", dark=True)
    add_two_column(
        s9,
        "Where Zariya Is Better",
        [
            "Clear articulation of quality and emotional identity",
            "Direct social messaging channels already integrated",
            "Simple, friction-light product discovery",
            "Clean baseline storefront without excessive clutter",
            "Ability to move quickly on Shopify stack",
        ],
        "Where Zariya Is Behind",
        [
            "Luxury brand proof and authority systems",
            "Editorial storytelling and style universe depth",
            "Premium conversion architecture on PDP",
            "Consistency in copy, policy quality, and trust polish",
            "Retention ecosystem beyond first purchase",
        ],
    )

    # Slide 10: 30-60-90 day plan
    s10 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s10, COLORS["bg_light"])
    add_header_bar(s10, "Action Plan: 30-60-90 Days", dark=True)
    add_bullets(
        s10,
        [
            "0-30 days: fix About page, rewrite policy copy, add trust badges, improve hero proposition",
            "0-30 days: redesign discount communication to avoid anti-luxury framing",
            "31-60 days: upgrade PDP with material, fit, care, and craftsmanship modules",
            "31-60 days: deploy UGC/reviews system with premium moderation and visual consistency",
            "61-90 days: launch drop calendar, referral loop, and founder-led content series",
            "61-90 days: set KPI dashboard and weekly growth-review ritual",
        ],
        dark=True,
    )

    # Slide 11: KPI dashboard
    s11 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s11, COLORS["bg_light"])
    add_header_bar(s11, "KPIs To Track", dark=True)
    add_two_column(
        s11,
        "Conversion Quality",
        [
            "PDP to cart rate",
            "Checkout completion rate",
            "Average order value",
            "Full-price sell-through vs discounted sell-through",
            "Return rate by SKU",
        ],
        "Brand Health",
        [
            "Repeat purchase rate (30/60/90 day)",
            "NPS or post-purchase satisfaction",
            "UGC volume and verified review velocity",
            "Direct traffic share and branded search trend",
            "Contribution margin after shipping and returns",
        ],
    )

    # Slide 12: Evidence appendix
    s12 = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(s12, COLORS["bg_light"])
    add_header_bar(s12, "Evidence Snapshot From Crawl", dark=True)
    add_bullets(
        s12,
        [
            "Homepage copy includes: Why Zariya, Quality Assurance, Brand Credibility, Safe Checkout",
            "Collection page shows repeated 41% OFF and price Rs. 699 vs Rs. 1,199 presentation",
            "Contact pathways available: Instagram, WhatsApp, email, phone",
            "Shipping policy present: process 1-3 business days, delivery 3-6 business days",
            "Refund policy present: 7-day return window, customer-paid return shipping unless issue at brand side",
            "About page crawl result: appears as 404 page, requiring immediate correction",
        ],
        dark=True,
    )

    prs.save(OUT_FILE)
    print(f"Created {OUT_FILE}")


if __name__ == "__main__":
    main()
